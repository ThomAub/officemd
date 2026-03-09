#!/usr/bin/env python3
"""Generate showcase OOXML files for examples and CLI demos.

Run from the OfficeMD workspace root:
    uv run --with python-docx --with openpyxl --with python-pptx python examples/generate_data.py
"""

from __future__ import annotations

from datetime import date, datetime
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET
import csv

from docx import Document
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches

PKG_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

REL_COMMENT_AUTHORS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors"
)
REL_COMMENTS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"

CT_COMMENT_AUTHORS = (
    "application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"
)
CT_COMMENTS = "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"


def _write_docx(path: Path) -> None:
    doc = Document()

    header = doc.sections[0].header
    header_paragraph = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    header_paragraph.text = "OfficeMD Showcase Header"

    doc.add_heading("Quarterly Operations Summary", level=1)
    doc.add_paragraph(
        "This fixture is used to demo DOCX extraction to IR and markdown."
    )

    comment_para = doc.add_paragraph(
        "Please validate this sentence before publishing the report."
    )
    doc.add_comment(
        comment_para.runs,
        text="Example DOCX comment captured as markdown footnote.",
        author="Reviewer",
        initials="RV",
    )

    table = doc.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Region"
    table.rows[0].cells[1].text = "Revenue"
    table.rows[0].cells[2].text = "Growth"
    table.rows[1].cells[0].text = "EMEA"
    table.rows[1].cells[1].text = "$1200000"
    table.rows[1].cells[2].text = "12%"
    table.rows[2].cells[0].text = "NA"
    table.rows[2].cells[1].text = "$980000"
    table.rows[2].cells[2].text = "8%"

    doc.add_paragraph(
        "End of document. The parser should emit a body section and a header section."
    )
    doc.save(path)


def _write_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"

    headers = ["Product", "BaseAmount", "Rate", "Total", "Notes"]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True)

    ws.append(["Widget", 1200, 0.15, "=B2*(1+C2)", "Primary SKU"])
    ws.append(["Gadget", 850, 0.10, "=B3*(1+C3)", "Secondary SKU"])
    ws.append(["Service", 600, 0.20, "=B4*(1+C4)", "Recurring"])

    ws["B2"].number_format = "$#,##0.00"
    ws["B3"].number_format = "$#,##0.00"
    ws["B4"].number_format = "$#,##0.00"
    ws["C2"].number_format = "0.00%"
    ws["C3"].number_format = "0.00%"
    ws["C4"].number_format = "0.00%"
    ws["D2"].number_format = "$#,##0.00"
    ws["D3"].number_format = "$#,##0.00"
    ws["D4"].number_format = "$#,##0.00"
    ws["A2"].comment = Comment("This row has a cell comment (Excel note).", "Analyst")

    ws["A6"] = "Project Wiki"
    ws["A6"].hyperlink = "https://example.com/wiki"
    ws["A6"].style = "Hyperlink"

    ws.freeze_panes = "A2"
    ws.column_dimensions["A"].width = 18
    ws.column_dimensions["B"].width = 14
    ws.column_dimensions["C"].width = 10
    ws.column_dimensions["D"].width = 14
    ws.column_dimensions["E"].width = 24

    ws2 = wb.create_sheet("Summary")
    ws2.append(["Metric", "Value"])
    ws2["A1"].font = Font(bold=True)
    ws2["B1"].font = Font(bold=True)
    ws2["A2"] = "ReportDate"
    ws2["B2"] = date(2026, 2, 1)
    ws2["B2"].number_format = "yyyy-mm-dd"
    ws2["A3"] = "RunAt"
    ws2["B3"] = datetime(2026, 2, 1, 14, 30, 0)
    ws2["B3"].number_format = "yyyy-mm-dd hh:mm"
    ws2["A4"] = "AverageRate"
    ws2["B4"] = "=AVERAGE(Sales!C2:C4)"
    ws2["B4"].number_format = "0.00%"

    wb.save(path)


def _write_csv(path: Path) -> None:
    rows = [
        ["Product", "BaseAmount", "Rate", "Total", "Notes"],
        ["Widget", "1200", "0.15", "=B2*(1+C2)", "Primary SKU"],
        ["Gadget", "850", "0.10", "=B3*(1+C3)", "Secondary SKU"],
        ["Service", "600", "0.20", "=B4*(1+C4)", "Recurring"],
    ]
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.writer(handle)
        writer.writerows(rows)


def _xml_to_bytes(root: ET.Element) -> bytes:
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _add_relationship(
    rels_bytes: bytes,
    rel_type: str,
    target: str,
    target_mode: str | None = None,
) -> bytes:
    ET.register_namespace("", PKG_RELS_NS)
    root = ET.fromstring(rels_bytes)

    max_idx = 0
    rel_tag = f"{{{PKG_RELS_NS}}}Relationship"
    for rel in root.findall(rel_tag):
        rel_id = rel.attrib.get("Id", "")
        if rel_id.startswith("rId") and rel_id[3:].isdigit():
            max_idx = max(max_idx, int(rel_id[3:]))

    next_id = f"rId{max_idx + 1}"
    attrs = {"Id": next_id, "Type": rel_type, "Target": target}
    if target_mode is not None:
        attrs["TargetMode"] = target_mode
    ET.SubElement(root, rel_tag, attrs)
    return _xml_to_bytes(root)


def _add_content_type_override(
    content_types_bytes: bytes,
    part_name: str,
    content_type: str,
) -> bytes:
    ET.register_namespace("", CONTENT_TYPES_NS)
    root = ET.fromstring(content_types_bytes)
    override_tag = f"{{{CONTENT_TYPES_NS}}}Override"

    for override in root.findall(override_tag):
        if override.attrib.get("PartName") == part_name:
            return _xml_to_bytes(root)

    ET.SubElement(
        root,
        override_tag,
        {"PartName": part_name, "ContentType": content_type},
    )
    return _xml_to_bytes(root)


def _write_pptx(path: Path) -> None:
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Revenue is up 12% quarter over quarter."

    table_shape = slide.shapes.add_table(3, 3, Inches(0.7), Inches(3.0), Inches(8.5), Inches(1.4))
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(0, 2).text = "Growth"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "$1.2M"
    table.cell(1, 2).text = "12%"
    table.cell(2, 0).text = "NA"
    table.cell(2, 1).text = "$0.98M"
    table.cell(2, 2).text = "8%"

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Speaker note: Call out that EMEA outperformed expectations."
    )

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Roadmap"
    textbox = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.2))
    textbox.text_frame.text = "Next: launch SKU refresh and expand partner channels."

    prs.save(path)

    with ZipFile(path, "r") as zin:
        parts = {name: zin.read(name) for name in zin.namelist()}

    presentation_rels = "ppt/_rels/presentation.xml.rels"
    slide1_rels = "ppt/slides/_rels/slide1.xml.rels"
    content_types = "[Content_Types].xml"

    parts[presentation_rels] = _add_relationship(
        parts[presentation_rels],
        REL_COMMENT_AUTHORS,
        "commentAuthors.xml",
    )
    parts[slide1_rels] = _add_relationship(
        parts[slide1_rels],
        REL_COMMENTS,
        "../comments/comment1.xml",
    )

    parts["ppt/commentAuthors.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8"?>\n'
        b'<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">\n'
        b'  <p:cmAuthor id="0" name="Presenter"/>\n'
        b"</p:cmAuthorLst>\n"
    )
    parts["ppt/comments/comment1.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8"?>\n'
        b'<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">\n'
        b'  <p:cm authorId="0"><p:text>Add one slide on operating margin.</p:text></p:cm>\n'
        b"</p:cmLst>\n"
    )

    parts[content_types] = _add_content_type_override(
        parts[content_types],
        "/ppt/commentAuthors.xml",
        CT_COMMENT_AUTHORS,
    )
    parts[content_types] = _add_content_type_override(
        parts[content_types],
        "/ppt/comments/comment1.xml",
        CT_COMMENTS,
    )

    with ZipFile(path, "w", ZIP_DEFLATED) as zout:
        for name in sorted(parts):
            zout.writestr(name, parts[name])


def main() -> None:
    output_dir = Path(__file__).parent / "data"
    output_dir.mkdir(parents=True, exist_ok=True)

    docx_path = output_dir / "showcase.docx"
    xlsx_path = output_dir / "showcase.xlsx"
    csv_path = output_dir / "showcase.csv"
    pptx_path = output_dir / "showcase.pptx"

    _write_docx(docx_path)
    _write_xlsx(xlsx_path)
    _write_csv(csv_path)
    _write_pptx(pptx_path)

    print(f"Generated {docx_path}")
    print(f"Generated {xlsx_path}")
    print(f"Generated {csv_path}")
    print(f"Generated {pptx_path}")


if __name__ == "__main__":
    main()
