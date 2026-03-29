# /// script
# requires-python = ">=3.12"
# dependencies = ["python-docx==1.1.2", "python-pptx==1.0.2"]
# ///
"""Generate synthetic benchmark corpus files.

Usage:
    uv run benchmark/generate-corpus.py [--output-dir benchmark/corpus]
"""

import argparse
import csv
import random
from pathlib import Path


def generate_customers_csv(output_dir: Path) -> None:
    """Generate a 1000-row customer CSV (~87KB)."""
    path = output_dir / "customers-1000.csv"
    if path.exists():
        print(f"  skip {path.name} (exists)")
        return

    random.seed(42)
    regions = ["North America", "Europe", "Asia Pacific", "Latin America", "Middle East"]
    products = ["Enterprise Plan", "Pro Plan", "Starter Plan", "API Access", "Support Add-on"]
    statuses = ["Active", "Churned", "Trial", "Paused"]

    with open(path, "w", newline="") as f:
        w = csv.writer(f)
        w.writerow(
            ["ID", "Company", "Region", "Product", "MRR", "Employees", "Status", "SignupDate", "LastLogin"]
        )
        for i in range(1, 1001):
            w.writerow([
                f"CUST-{i:04d}",
                f"Company {i}",
                random.choice(regions),
                random.choice(products),
                round(random.uniform(99, 9999), 2),
                random.randint(5, 10000),
                random.choice(statuses),
                f"2023-{random.randint(1, 12):02d}-{random.randint(1, 28):02d}",
                f"2024-{random.randint(1, 12):02d}-{random.randint(1, 28):02d}",
            ])

    print(f"  wrote {path.name}")


def generate_tech_architecture_docx(output_dir: Path) -> None:
    """Generate a multi-section technical document (~37KB)."""
    path = output_dir / "tech-architecture.docx"
    if path.exists():
        print(f"  skip {path.name} (exists)")
        return

    from docx import Document

    doc = Document()
    doc.add_heading("Technical Architecture Document", 0)
    doc.add_paragraph("Version 2.1 | Last Updated: December 2024")

    doc.add_heading("1. System Overview", level=1)
    doc.add_paragraph(
        "The platform consists of a microservices architecture deployed on Kubernetes, "
        "with a React frontend communicating via GraphQL APIs. The system handles "
        "approximately 50,000 requests per second at peak load, serving customers "
        "across 12 geographic regions."
    )

    doc.add_heading("1.1 Core Components", level=2)
    doc.add_paragraph("The following table describes the main system components:")
    table = doc.add_table(rows=8, cols=4)
    table.style = "Table Grid"
    for i, h in enumerate(["Component", "Technology", "Instances", "Purpose"]):
        table.cell(0, i).text = h
    for i, row in enumerate([
        ["API Gateway", "Kong + Envoy", "6", "Request routing, rate limiting, auth"],
        ["Auth Service", "Go + JWT", "4", "Authentication and authorization"],
        ["Core API", "Node.js + TypeScript", "12", "Business logic and data access"],
        ["Search Engine", "Elasticsearch 8.x", "3", "Full-text search and analytics"],
        ["Message Queue", "Apache Kafka", "5", "Event streaming and async processing"],
        ["Cache Layer", "Redis Cluster", "6", "Session storage, caching, rate limiting"],
        ["Database", "PostgreSQL 16", "3", "Primary data store with read replicas"],
    ]):
        for j, val in enumerate(row):
            table.cell(i + 1, j).text = val

    doc.add_heading("2. Performance Requirements", level=1)
    doc.add_paragraph("The system must meet the following SLAs:")
    for item in [
        "API response time: p50 < 50ms, p99 < 200ms",
        "Availability: 99.99% uptime (< 52 minutes downtime/year)",
        "Throughput: 50,000 RPS sustained, 100,000 RPS burst",
        "Data durability: 99.999999999% (11 nines)",
        "Recovery Point Objective (RPO): < 1 minute",
        "Recovery Time Objective (RTO): < 5 minutes",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_heading("3. Data Model", level=1)
    doc.add_paragraph("The primary database schema consists of the following entities:")
    table2 = doc.add_table(rows=7, cols=3)
    table2.style = "Table Grid"
    for i, h in enumerate(["Entity", "Rows (approx)", "Key Relations"]):
        table2.cell(0, i).text = h
    for i, row in enumerate([
        ["Users", "3.8M", "Has many Orders, Sessions, Preferences"],
        ["Orders", "47M", "Belongs to User, Has many LineItems"],
        ["Products", "125K", "Has many Variants, Categories, Reviews"],
        ["Reviews", "12M", "Belongs to User and Product"],
        ["Sessions", "890M", "Belongs to User, Has many Events"],
        ["Events", "4.2B", "Belongs to Session, polymorphic type"],
    ]):
        for j, val in enumerate(row):
            table2.cell(i + 1, j).text = val

    doc.add_heading("4. Deployment Architecture", level=1)
    doc.add_paragraph(
        "All services are containerized and deployed via Helm charts to AWS EKS clusters. "
        "We maintain three environments: development, staging, and production."
    )
    table3 = doc.add_table(rows=6, cols=3)
    table3.style = "Table Grid"
    for i, h in enumerate(["Service", "Monthly Cost", "Notes"]):
        table3.cell(0, i).text = h
    for i, row in enumerate([
        ["AWS EKS", "$12,400", "3 clusters, 45 nodes total"],
        ["RDS PostgreSQL", "$8,200", "Multi-AZ, 3 read replicas"],
        ["ElastiCache Redis", "$3,100", "6-node cluster"],
        ["CloudFront CDN", "$2,800", "~15TB/month transfer"],
        ["Total", "$26,500", "Excluding data transfer and support"],
    ]):
        for j, val in enumerate(row):
            table3.cell(i + 1, j).text = val

    doc.save(str(path))
    print(f"  wrote {path.name}")


def generate_q4_business_review_pptx(output_dir: Path) -> None:
    """Generate a Q4 business review presentation (~34KB)."""
    path = output_dir / "q4-business-review.pptx"
    if path.exists():
        print(f"  skip {path.name} (exists)")
        return

    from pptx import Presentation

    prs = Presentation()
    slides_data = [
        ("Q4 2024 Business Review", "Annual Performance Summary\nPrepared by Analytics Team"),
        ("Agenda", "1. Revenue Overview\n2. Customer Metrics\n3. Product Performance\n4. Market Expansion\n5. 2025 Roadmap"),
        ("Revenue Overview", "Total Revenue: $45.2M (+23% YoY)\nRecurring Revenue: $38.1M (+31% YoY)\nNew Business: $7.1M\nNet Revenue Retention: 118%\nGross Margin: 72.4%"),
        ("Customer Metrics", "Total Customers: 3,847 (+42% YoY)\nEnterprise: 312 accounts\nMid-Market: 1,205 accounts\nSMB: 2,330 accounts\nChurn Rate: 2.1% (down from 3.8%)\nNPS Score: 72"),
        ("Product Performance", "Core Platform: 89% adoption\nAPI Usage: 2.3B calls/month (+156%)\nMobile App: 45K DAU\nNew Features Shipped: 47\nUptime: 99.97%\nAvg Response Time: 142ms"),
        ("Market Expansion", "Launched in 3 new markets: Germany, Japan, Brazil\nEMEA Revenue: $8.2M (+67%)\nAPAC Revenue: $4.1M (+89%)\nPartner Channel: 23 new integrations\nSOC 2 Type II Certified"),
        ("2025 Roadmap", "Q1: AI-powered analytics engine\nQ2: Enterprise SSO & SCIM\nQ3: Real-time collaboration\nQ4: On-premise deployment option\n\nTarget: $65M ARR by EOY 2025"),
        ("Team & Hiring", "Current Headcount: 187\nEngineering: 89\nSales: 34\nCustomer Success: 28\nG&A: 36\n\n2025 Hiring Plan: +65 heads"),
    ]
    for title, body in slides_data:
        layout = prs.slide_layouts[0] if title == slides_data[0][0] else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    prs.save(str(path))
    print(f"  wrote {path.name}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate synthetic benchmark corpus files.")
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path(__file__).resolve().parent / "corpus",
        help="Directory to write generated files (default: benchmark/corpus)",
    )
    args = parser.parse_args()
    args.output_dir.mkdir(parents=True, exist_ok=True)

    print("Generating synthetic files:")
    generate_customers_csv(args.output_dir)
    generate_tech_architecture_docx(args.output_dir)
    generate_q4_business_review_pptx(args.output_dir)


if __name__ == "__main__":
    main()
